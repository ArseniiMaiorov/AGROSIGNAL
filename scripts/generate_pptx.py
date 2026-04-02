import os
import subprocess
import sys

# Пытаемся импортировать python-pptx, если нет - устанавливаем
try:
    import pptx
except ImportError:
    print("Устанавливаем библиотеку python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Стилизация текста
    for shape in [title, subtitle]:
        shape.text_frame.paragraphs[0].font.name = 'Arial'
        if shape == title:
            shape.text_frame.paragraphs[0].font.size = Pt(40)
            shape.text_frame.paragraphs[0].font.bold = True
        else:
            shape.text_frame.paragraphs[0].font.size = Pt(24)
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

def add_image_slide(prs, title_text, img_path, caption_text=""):
    slide_layout = prs.slide_layouts[5] # Blank with Title
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.bold = True
    
    if os.path.exists(img_path):
        # Добавляем изображение, подгоняя под размер слайда
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5.0)
        width = Inches(9.0)
        
        try:
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
            # Выравниваем по центру по горизонтали
            pic.left = int((prs.slide_width - pic.width) / 2)
            
            # Если есть подпись - добавляем под картинку
            if caption_text:
                txBox = slide.shapes.add_textbox(Inches(0.5), top + pic.height + Inches(0.1), Inches(9.0), Inches(0.8))
                tf = txBox.text_frame
                tf.text = caption_text
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(16)
                p.font.name = 'Arial'
                p.font.color.rgb = RGBColor(80, 80, 80)
        except Exception as e:
            print(f"Ошибка при добавлении {img_path}: {e}")
    else:
        # Если картинки нет, просто пишем текст
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        tf = txBox.text_frame
        tf.text = f"(Изображение не найдено: {img_path})\n{caption_text}"

def create_presentation():
    # Путь к изображениям
    # Подразумеваем, что скрипт запускается из Диплом/AGROSIGNAL
    base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', 'ПЗ', 'figures'))
    
    prs = Presentation()
    # Устанавливаем размер 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 1. Титульный слайд
    add_title_slide(prs, 
                   "ТерраINFO: Информационная поддержка агрономических решений", 
                   "Основано на спутниковых, погодных и пользовательских данных\nВКР")
                   
    # 2. Проблематика
    add_image_slide(prs, "Проблематика и архитектура системы", 
                   os.path.join(base_dir, 'system_architecture.png'),
                   "Единая экосистема для устранения разрозненности данных и автоматизации аналитики")
                   
    # 3. Десктоп-оболочка
    add_image_slide(prs, "Полноэкранный агрономический дашборд", 
                   os.path.join(base_dir, 'desktop_shell_annotated.png'),
                   "Рабочее пространство и единый интерфейс пользователя")
                   
    # 4. Детект и Картография
    add_image_slide(prs, "Спутниковая картография и геоаналитика", 
                   os.path.join(base_dir, 'map_layers.png'),
                   "Мультиспектральные слои и индексы вегетации")
                   
    # 5. Панель управления (контроль слоев)
    add_image_slide(prs, "Определение границ: Панель управления", 
                   os.path.join(base_dir, 'control_panel_annotated.png'),
                   "Настройка автоматического оконтуривания полей и слоев")
                   
    # 6. Алгоритм сегментации и признаки
    add_image_slide(prs, "Обучающие признаки для алгоритма сегментации", 
                   os.path.join(base_dir, 'perm_feature_maps.png'),
                   "Формирование многоканальных тензоров для модели BoundaryUNet")
                   
    # 7. Результаты обучения модели
    add_image_slide(prs, "BoundaryUNet v3: Результаты обучения", 
                   os.path.join(base_dir, 'metrics_xy_annotated.png'),
                   "Существенный прирост метрик: +64% F1-score, снижение ошибки HD95 на 59%")
                   
    # 8. Панель поля (Аналитика)
    add_image_slide(prs, "Детальная аналитика растительности поля", 
                   os.path.join(base_dir, 'field_overview_annotated.png'),
                   "Анализ индексных кривых и фаз развития культур")
                   
    # 9. Прогноз и мониторинг
    add_image_slide(prs, "Мониторинг погоды и климатических рисков", 
                   os.path.join(base_dir, 'weather_panel_annotated.png'),
                   "Ежедневная динамика, ET0, GDD и стрессовые индексы")
                   
    # 10. Прогнозирование урожайности
    add_image_slide(prs, "Предиктивный прогноз урожайности", 
                   os.path.join(base_dir, 'forecast_panel_annotated.png'),
                   "Расчёт потенциала на основе накопленных данных")
                   
    # 11. Сценарное моделирование
    add_image_slide(prs, "Сценарное моделирование (Что-если?)", 
                   os.path.join(base_dir, 'scenario_panel_annotated.png'),
                   "Моделирование влияния агрономических мероприятий на итоговый сбор")
                   
    # 12. Архив мероприятий
    add_image_slide(prs, "Журнал мероприятий и событий", 
                   os.path.join(base_dir, 'archive_panel_annotated.png'),
                   "Учёт выполненных работ и внесения удобрений")

    # 13. Заключение
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Заключение и Эффект"
    content = slide.placeholders[1]
    content.text = ("• Полный цикл: детект границ, погодная аналитика, сценарный прогноз\n"
                    "• Региональная адаптация ИИ завершена\n"
                    "• Экономический эффект: прирост сбора ~2%\n"
                    "• Окупаемость внедрения: ~7 месяцев\n\n"
                    "Система ТерраINFO полностью готова к защите и эксплуатации")

    out_file = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', 'ПЗ', 'TerraINFO_Presentation.pptx'))
    prs.save(out_file)
    print(f"Презентация успешно создана: {out_file}")

if __name__ == "__main__":
    create_presentation()
